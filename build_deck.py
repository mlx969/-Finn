# -*- coding: utf-8 -*-
"""
Day 10 (2026-08-12, MVP 冲刺) — 生成「乐天知性(安心答)」参赛方案 deck + PDF。

生成物：
  - presentation.pptx  : 18 页浓缩方案（12 章报告 -> 评委可快速读完）
  - proposal.pdf       : 1 页执行摘要（PDF 备选）

说明：
  - 评测/红队真题图因 A5(kb_adult 检索空)未闭环，仍用 report_assets/ 的「模拟占位」图；
    图下保留「模拟数据·待真值回填」标注。A5 闭环后按 REPORT_DATA_SLOTS.md 一键替换。
  - 全文原创，不复制任何创作者原话；严守 18+ 硬门槛 / 七条安全规则红线（仅说明，不演示绕过）。
  - 不写任何 API Key（本脚本无 Key 需求）。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.abspath(__file__))
ASSET_DIR = os.path.join(ROOT, "report_assets")

# 色板（与 index.html 暖珊瑚奶油一致）
CORAL   = RGBColor(0xF2, 0x52, 0x7A)   # 主品牌色
CORAL2  = RGBColor(0xFF, 0x6B, 0x8A)
CREAM   = RGBColor(0xFF, 0xF5, 0xF8)
INK     = RGBColor(0x2E, 0x2A, 0x33)
PLUM    = RGBColor(0x7A, 0x4A, 0x57)
GREY    = RGBColor(0x93, 0x8B, 0x97)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
REDCARD = RGBColor(0xB0, 0x30, 0x3A)   # 红线强调

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def _set(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT,
         font=FONT, space_after=6, level=0, line=1.15):
    p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    p.line_spacing = line
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return p


def fill_bg(s, color=WHITE):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def title_bar(s, kicker, title, num):
    rect(s, 0, 0, SW, Inches(1.15), CREAM)
    rect(s, 0, 0, Inches(0.16), Inches(1.15), CORAL)
    tb = s.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(11.6), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, kicker, 12, CORAL, bold=True, space_after=2)
    _set(tf, title, 25, INK, bold=True, space_after=0)
    # 页码
    pg = s.shapes.add_textbox(Inches(12.4), Inches(0.35), Inches(0.8), Inches(0.6))
    _set(pg.text_frame, num, 13, GREY, bold=True, align=PP_ALIGN.RIGHT)


def bullets(s, items, x, y, w, h, size=15, gap=8, color=INK, header=None):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True
    first = True
    if header:
        _set(tf, header, size + 2, PLUM, bold=True, space_after=gap + 2)
        first = False
    for it in items:
        if isinstance(it, tuple):
            text, lvl = it
        else:
            text, lvl = it, 0
        p = tf.paragraphs[0] if (first and len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        p.line_spacing = 1.12
        bullet = ("● " if lvl == 0 else "– ")
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return box


def footer(s, note=""):
    f = s.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(12.4), Inches(0.35))
    _set(f.text_frame, "乐天知性 · 安心答  |  GOAI 世界人工智能开源大赛 · 赛道二 无界应用  |  MIT 开源" + (("  |  " + note) if note else ""),
         9, GREY, space_after=0)


def placeholder_note(s, x, y, w, text="⚠️ 模拟数据·待真值回填（A5 闭环后替换）"):
    box = s.shapes.add_textbox(x, y, w, Inches(0.4))
    _set(box.text_frame, text, 10, REDCARD, bold=True, space_after=0)


# ============================ SLIDES ============================

# 1. 封面
s = slide(); fill_bg(s, CREAM)
rect(s, 0, 0, SW, Inches(0.22), CORAL)
rect(s, 0, Inches(7.28), SW, Inches(0.22), CORAL)
t = s.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.6))
tf = t.text_frame; tf.word_wrap = True
_set(tf, "乐天知性 · 安心答", 46, INK, bold=True, space_after=6)
_set(tf, "AI + 性健康与亲密关系 · 社会学升维陪伴 + 适龄双版本", 20, PLUM, space_after=14)
_set(tf, "GOAI 世界人工智能开源大赛  ·  赛道二 无界应用（Boundless Agents）", 16, GREY)
sub = s.shapes.add_textbox(Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.2))
tf2 = sub.text_frame; tf2.word_wrap = True
_set(tf2, "Slogan：乐天知性，安心作答 —— 懂身体，也懂你", 15, CORAL2, bold=True, space_after=4)
_set(tf2, "参赛者：个人参赛（3073812933@qq.com）  ·  开源协议：MIT  ·  版本 v2.0", 13, GREY)
footer(s)

# 2. 摘要 / 一句话定位
s = slide(); fill_bg(s)
title_bar(s, "EXECUTIVE SUMMARY", "我们解决什么 / 怎么做", "02")
bullets(s, [
    "「乐天知性」是面向性健康与亲密关系场景的 AI 助手（产品名「安心答」），核心差异在于引入社会学升维视角。",
    "把个体的情感与性困惑，放回原生家庭、阶层处境、身份位置等结构性变量中理解——而非停留于生理知识或猎奇解答。",
    "双版本架构：18+ 版（升维对话 + 小众文化科普 + 18+ 论坛 + 成人展资讯）与未成年版（适龄科普 + 动漫/游戏化，零成人内容）。",
    "全球首个将「升维视角 + 适龄双版本 + 社区 + 审核」整合的开源尝试；中文性教育/亲密关系垂类在 GitHub/HF 近乎空白。",
    "可体验（Demo）· 可复制（开源）· 可落地（真实场景价值）；所有入口强制 18+ 年龄认证，严守 7 条安全规则。",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=16, gap=12)
footer(s)

# 3. 问题定义
s = slide(); fill_bg(s)
title_bar(s, "01 · 问题定义", "长期缺位 × 视角空白 × 适龄段空白", "03")
bullets(s, [
    "长期缺位：性教育长期「学校不敢讲、家庭不会讲、网络乱讲」，权威覆盖率极低，意外妊娠/STI/情感创伤风险却高。",
    "市面两类极端：一端教科书式生理罗列（枯燥不解决问题）；另一端猎奇/擦边（无护栏、夹带 PUA/物化框架）。",
    "视角空白：几乎所有内容只在「知识」层打转，缺把困惑与原生家庭/阶层/身份挂钩的「升维」解释——而这正是年轻人最想被理解的维度。",
    ("洞察：青少年普遍早熟但概念最模糊，处于「好奇+追求刺激+不懂后果」高风险区；现有轮子几乎全 18+ 成人向，适龄段全球空白。", 0),
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=15, gap=11)
footer(s)

# 4. 生态对标
s = slide(); fill_bg(s)
title_bar(s, "02 · 生态对标", "红海里的空白（GitHub / HuggingFace 调研）", "04")
bullets(s, [
    "心理/情感垂类已是红海：EmoLLM、SoulChat2.0(ACL2025)、MeChat/SMILE、CPsyCoun 等——可直接借鉴工程范式而非重复造轮子。",
    "性教育垂类全球空白：开源「性健康 chatbot」仅 5 个、全在海外、且全部只做生理问答。",
    "竞品共性缺失：SafeBubble(丹麦,隐私即架构) / TroyHealthLink(美,风险自评) / Aarogya_Mitram(印,8 Agent) … 都无「升维+社区+适龄版」。",
    "结论：技术水位≈硕士毕设，Dify/FastGPT 低代码 3–7 天可复刻；本项目定位为全球首个「升维+双版本+社区+审核」四者整合的开源。",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=15, gap=11)
footer(s)

# 5. 差异化定位 — 第四类
s = slide(); fill_bg(s)
title_bar(s, "03 · 差异化定位 ★", "第四类内容 + 五层架构", "05")
bullets(s, [
    "市面三类我们都不做：教科书式（枯燥）✗ / 猎奇擦边（无护栏、夹带 PUA）✗ / PUA物化（教算计对方）✗。",
    "本项目做第四类：知识 + 结构视角 + 适龄分级 + 安全 + 社区。",
    "五层架构：①知识层(RAG权威资料) ②升维层(情感社会学框架) ③多元层(SM/性少数去污名,consent前置) ④社群层(18+论坛+展讯) ⑤适龄层(未成年版动漫/游戏化,物理隔离)。",
    "升维视角定义(原创)：把「你现在的情感模式」视为成长环境与结构性处境的产物，而非个人缺陷；给用户更宽的理解框架，促其自主决策。",
    "小众文化定位为「性与亲密关系安全教育」：consent/安全/沟通/健康风险/去污名——是科普不是猎奇(仅限18+)。",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=14, gap=9)
footer(s)

# 6. 双版本壁垒
s = slide(); fill_bg(s)
title_bar(s, "03.4 · 核心壁垒 ★", "双版本：竞品零覆盖的差异化空白", "06")
bullets(s, [
    "洞察：早熟 ≠ 懂。15 岁即外表成熟，但性知识概念最模糊，处于高风险区；初中即有因无知好奇偷尝禁果的现象。",
    "机会：GitHub/HF 的性健康 bot 几乎全 18+ 成人向，适龄段(尤其12–18)无人占；把市面性教育课程用动漫风+游戏化重构，全球稀缺。",
    "实现：身份/年龄认证后分流，未成年版与成年版后端双应用物理隔离（详见第9章）。",
    "红线(最高优先级，不可逾越)：绝不性化未成年人；未成年版零色情/零性暗示/零SM；强制年龄门+监护人提示+数据最小化。",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=15, gap=12)
footer(s)

# 7. 系统架构
s = slide(); fill_bg(s)
title_bar(s, "04 · 系统架构", "单系统 · 双隔离子应用（纵深防御）", "07")
bullets(s, [
    "统一品牌入口(index.html / Dify WebApp) → 身份/年龄认证(第一道墙) → 分流未成年(<18)/成年(≥18)。",
    "未成年版：青少前端(动漫/游戏化科普,零成人) + Dify应用A(kb_minor仅适龄,零SM/零性暗示)。",
    "成年版：成年前端(升维对话/SM/论坛/展讯) + Dify应用B(kb_adult全量,含SM科普/多元/论坛/展讯)。",
    "纵深防御：前端判断+后端隔离两层墙。未成年版应用根本不挂 kb_adult、不配成人节点——即使年龄门被绕过也泄露不了成人内容。",
    "选型 Chatflow(非Agent)：涉性+未成年必须确定性安全拦截，节点流程显式可控、风险转介可做独立节点、评测输出稳定。",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=14, gap=9)
footer(s)

# 8. 成年版链路
s = slide(); fill_bg(s)
title_bar(s, "04.2 · 已跑通 ✅", "成年版 Chatflow 实际链路（5 节点）", "08")
bullets(s, [
    "[开始/用户输入] → [知识检索 kb_adult, TopK=3, Score≥0.6] → [条件分支 IF检索有果 / ELSE兜底]。",
    "IF 分支 → [LLM: 阶跃 Step 3.7 Flash + SYSTEM_PROMPT(7条安全规则+升维)] → [直接回复(含来源标注)]。",
    "现状：应用创建、LLM挂SYSTEM_PROMPT、阶跃绑定、知识检索+条件分支 5 节点连线均已完成；基础对话实测跑通(升维风格到位)。",
    "技术栈：阶跃 StepFun(step-router-v1/Step3.7 Flash) + Dify Chatflow(主力,可私有化) + FastGPT(对照) + ShieldLM-6B+LlamaGuard S1–S13+DFA(四层审核) + 单文件index.html + Cloudflare Pages(部署)。",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=14, gap=10)
footer(s)

# 9. 知识库
s = slide(); fill_bg(s)
title_bar(s, "05 · 知识库构建", "权威来源 × 双库物理隔离 × 原创重写", "09")
bullets(s, [
    "权威来源：WHO 性健康与生殖健康指南、中国计生协/疾控公开科普、多元性文化去污名化学术综述。",
    "双库物理隔离：kb_adult(WHO/计生协/疾控+多元+升维,18+版,已在Dify创建) / kb_minor(刘文利读本/UNESCO ITGSE/保护豆豆/LoveMatters青春版/青爱工程,未成年版,规划中,先12–15)。",
    "数据策略(引轮子+现状润色)：只抽 PsyQA/CPsyCoun/SoulChat/EmoLLM 等骨架，用升维视角+2025–26现状重写为原创问答(our_corpus,随MIT开源)。",
    "RAG 关键设计：切分带元数据(来源/主题/适用年龄/版权) + 生成强制引用降幻觉 + 非本领域触发意图闸门拒答。",
    "人卫教材仅索引/摘要引用，禁止全文灌库（版权边界，见 DATA_STRATEGY.md）。",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=14, gap=9)
footer(s)

# 10. 安全与合规（核心）
s = slide(); fill_bg(s)
title_bar(s, "06 · 安全与合规 ★", "18+ 硬门槛 + 七条强制规则 + 四层审核", "10")
bullets(s, [
    "身份/年龄认证(第一道墙)：Demo 自声明+生日分级；上线预留实名接口；未成年版不渲染任何18+入口（防绕过）。",
    "七条强制安全规则(SYSTEM_PROMPT.md)：①AGE_GATE 仅服务18+ ②NO_PORN 绝不露骨 ③CONSENT_FIRST 知情同意 ④SELF_HARM_REDIRECT ⑤VIOLENCE_REDIRECT ⑥MEDICAL_REDIRECT ⑦SOURCE_TAG 标注来源。",
    "四层审核：L1 DFA敏感词(6万,毫秒拦截) / L2 ShieldLM-6B(语义二审) / L3 人工抽审+危机升级Agent / L4 用户信用+举报。",
    "内容分级采用 Llama Guard MLCommons S1–S13：S4(儿童性剥削)零容忍硬拦，S3/S10/S11 全覆盖高危场景。",
    "危机转介：检测自伤/家暴/医疗紧急 → 输出专业资源指引，不替代专业干预（热线号码已建白名单并设过期复核）。",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=14, gap=9)
footer(s, "红线③：API Key 仅内存、不落库")

# 11. 社区与论坛
s = slide(); fill_bg(s)
title_bar(s, "07 · 社区与论坛（18+版）", "四版块 + 四层审核 + 成人展聚合", "11")
bullets(s, [
    "四个版块：①升维讨论(情感模式与原生家庭/阶层) ②安全教育(SM/多元文化科普) ③成人展资讯(公开活动聚合) ④互助答疑(经验分享)。",
    "门槛：18+ 强制 + 四层审核；小众文化区 NSFW 折叠 + 二次年龄验证。",
    "成人展模块只做公开活动资讯聚合（已调研青岛 + 2026–27 南方为主版图：上海API EXPO/ADC-EXPO、广州AAE、深圳SZI等）。",
    "边界：不做私域引流、不代办票务、以官方为准；合规说明见 ADULT_EXPO_COMPLIANCE.md。",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=15, gap=12)
footer(s)

# 12. 评测方案
s = slide(); fill_bg(s)
title_bar(s, "08 · 评测方案 ★", "三组对照 + 152 测试 + 47 红队 + 五维评分", "12")
bullets(s, [
    "三组对照：A 裸大模型(基线) / B Dify版(主方案,已实跑链路) / C FastGPT版(对照)。",
    "测试集 152 条(已生成)：知识37 / 升维37 / 安全38 / 拒答·转介40，四类均衡。",
    "红队 47 条(已生成)：越狱/色情诱导/伪装未成年/违法诱导/指令注入/PUA/对抗绕审/概念扭曲，每条标攻击类型与预期拦截层。",
    "五维评分 Rubric(0–2)：知识准确 / 升维质量 / 安全合规 / 拒答恰当 / 引用规范；红队以拦截率为核心指标。",
], Inches(0.5), Inches(1.4), Inches(7.6), Inches(5.0), size=14, gap=10)
# 右侧嵌入占位图
img_path = os.path.join(ASSET_DIR, "chart_5dim_模拟占位.png")
if os.path.exists(img_path):
    pic = s.shapes.add_picture(img_path, Inches(8.3), Inches(1.55), width=Inches(4.6))
    placeholder_note(s, Inches(8.3), Inches(5.7), Inches(4.6))
footer(s)

# 13. 评测图表（占位，整页）
s = slide(); fill_bg(s)
title_bar(s, "08 · 评测可视化（占位）", "五维得分对照 + 安全拦截率（模拟数据）", "13")
if os.path.exists(img_path):
    s.shapes.add_picture(img_path, Inches(0.5), Inches(1.35), width=Inches(7.4))
    placeholder_note(s, Inches(0.5), Inches(6.5), Inches(7.4))
bullets(s, [
    ("A5(kb_adult 检索空)未闭环 → 真值无法回填，本图与拦截率图均为模拟占位。", 0),
    "A5 闭环当天：按 REPORT_DATA_SLOTS.md 跑 eval_runner(152+47) → eval_scorer 出真值图与记分卡 → 本页整页替换为实测 v1。",
    "红队拦截率图(chart_intercept_模拟占位)因仅 SVG 格式，PPT 以本页五维图示意，PDF/报告引用 SVG 原件。",
    "评分器声明：本分为正则启发式自动初筛分，非人工评分，引用须注明。",
], Inches(8.1), Inches(1.5), Inches(4.9), Inches(5.0), size=13, gap=10, color=PLUM)
footer(s, "图下须保留「模拟数据·待真值回填」")

# 14. 双版本详解
s = slide(); fill_bg(s)
title_bar(s, "09 · 双版本架构详解", "适龄分级 + 动漫化 + 游戏化 + 实现", "14")
bullets(s, [
    "适龄分级：12–15(基础生理/青春期/身体权与边界/网络自护) / 15–18(进阶关系决策/consent/避孕防护基础/情感自尊)。",
    "6 大知识域(去性化、教育向)：生理发育与卫生 / 身体自主权与边界 / 同意与尊重 / 网络安全自护(裸聊勒索、隔空猥亵防范) / 情绪与关系 / 风险认知。",
    "动漫化(科普番剧)：参考《工作细胞》拟人化，把卵子/精子/月经/激素拟人讲生理，3–5分钟/集，视觉清新卡通，绝不走性感/萝莉风。",
    "游戏化(小课堂)：互动剧情选择 + Quiz 闯关 + 角色养成(无性化)；MVP 先出 1 个可玩互动页验证吸收效果。",
    "技术实现：Dify 双应用(应用A成年版已实跑 / 应用B青少版新建SYSTEM_PROMPT_MINOR) + index.html 认证路由+双版UI切换。",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=14, gap=9)
footer(s)

# 15. MVP 取舍
s = slide(); fill_bg(s)
title_bar(s, "09.6 · MVP 取舍（已采纳）", "设计完整 + 实现单版", "15")
bullets(s, [
    "MVP 阶段做「设计完整 + 实现单版」：报告与演示把双版本架构讲透(加分战略叙事)。",
    "成年版做深做实：152 条评测 + 47 条红队实测(待 A5 闭环跑分)。",
    "未成年版做「一个可点击入口 + 2–3 屏样例 + 1 个动漫风互动 demo 页」证明可行，不铺开第二套知识库与评测。",
    "理由：评委看的是深度+证据链；两套都半成品，不如一套有硬数据来得扎实。",
    "评测分层：eval_testset_minor.csv(适龄命中/边界提示/拒答成人内容) + 红队新增「伪装成年骗门」「未成年请求SM/色情」必须拒答。",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=15, gap=11)
footer(s)

# 16. 开源贡献 + 落地路径
s = slide(); fill_bg(s)
title_bar(s, "10–11 · 开源贡献 & 落地", "MIT 开源 · 比赛路径 · 真实上线分级", "16")
bullets(s, [
    "开源内容(MIT)：index.html 原型 / SYSTEM_PROMPT.md(+MINOR) / KNOWLEDGE_BASE.md+kb_index_v0.json / 评测红队套件 / 双版本与品牌设计文档。",
    "开源价值：公开升维视角提示词框架 + 四层审核架构 + 评测/红队集 + 双版本适龄分级范式，降低领域安全研究门槛。",
    "比赛路径：goaihz.com 报名 → 初赛(8/16截止) → 复赛(8/25–9/3) → 决赛(9/22–23 杭州 GOAI DAY)；配合 Datawhale 夏令营 baseline。",
    "真实上线路径(合规分级)：Demo级(本仓库+私有Dify)；上线级需 ICP+网文证+内容审查，部署建议 Cloudflare Pages，Key 经密钥管理服务注入不落库。",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=14, gap=10)
footer(s)

# 17. 风险与边界
s = slide(); fill_bg(s)
title_bar(s, "12 · 风险与边界", "如实标注，不夸大落地能力", "17")
bullets(s, [
    "内容合规(未成年人/色情)：高 → 18+硬门槛+四层审核+S4零容忍+双版本物理隔离。",
    "热线号码失效：高(仅危机转介处) → C4 已建白名单并设过期复核，核验前只给渠道类型+110/120。",
    "知识版权(教材全文)：中 → 仅索引/摘要引用。  医疗误答：中 → 强制转介+标注「非专业诊断」。",
    "数据溯源不足：中 → 附录 C1 待用户回溯原始出处(提交前必做)。  双版本排期：中 → MVP采用「设计完整+实现单版」缩减范围。",
    "边界声明：本项目是教育辅助工具，非医疗/心理咨询替代品；所有建议须结合专业资源。",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), size=14, gap=9)
footer(s)

# 18. 结尾 / 致评委
s = slide(); fill_bg(s, CREAM)
rect(s, 0, 0, SW, Inches(0.22), CORAL)
rect(s, 0, Inches(7.28), SW, Inches(0.22), CORAL)
t = s.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(3.2))
tf = t.text_frame; tf.word_wrap = True
_set(tf, "乐天知性，安心作答", 40, INK, bold=True, space_after=14)
_set(tf, "一个安全、可信、有社会学升维视角、适龄分级、带社区与活动资讯的性健康 AI 助手。", 18, PLUM, space_after=12)
_set(tf, "可体验 · 可复制(MIT 开源) · 可落地 —— 全球首个「升维+双版本+社区+审核」整合开源尝试。", 15, GREY, space_after=10)
_set(tf, "关键节点：🟢 8/15 MVP 冻结  ｜  🔴 8/16 初赛提交(goaihz.com)", 14, CORAL, bold=True)
footer(s)

# ============================ SAVE PPTX ============================
pptx_path = os.path.join(ROOT, "presentation.pptx")
prs.save(pptx_path)
print("SAVED PPTX:", pptx_path, "slides:", len(prs.slides._sldIdLst))
